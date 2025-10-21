from pathlib import Path
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from typing import Literal
from langchain_core.tools import tool
import re

@tool
def generate_report(
    file_type: Literal["pdf", "pptx"],
    title: str,
    sections: list[dict],
    output_path: str = "output_file"
) -> str:
    """
    Generate a visually appealing PDF or PPTX file with custom styles and layouts.
    The report will be saved inside a 'report' folder located in the same directory
    as this script.
    """

    # === Ensure the 'report' folder exists ===
    base_dir = Path(__file__).parent
    report_dir = base_dir / "report"
    report_dir.mkdir(exist_ok=True)

    # Full output file path (inside 'report' folder)
    file_path = report_dir / f"{output_path}.{file_type}"

    # === PDF GENERATION ===
    if file_type == "pdf":
        doc = SimpleDocTemplate(
            str(file_path),
            pagesize=A4,
            topMargin=50,
            bottomMargin=50,
            leftMargin=60,
            rightMargin=60,
        )

        styles = getSampleStyleSheet()

        title_style = ParagraphStyle(
            "TitleCustom",
            parent=styles["Title"],
            fontSize=22,
            leading=28,
            textColor=colors.HexColor("#2C3E50"),
            spaceAfter=20,
        )

        heading_style = ParagraphStyle(
            "HeadingCustom",
            parent=styles["Heading2"],
            fontSize=14,
            textColor=colors.HexColor("#1ABC9C"),
            spaceBefore=15,
            spaceAfter=8,
        )

        body_style = ParagraphStyle(
            "BodyCustom",
            parent=styles["BodyText"],
            fontSize=11,
            leading=16,
            textColor=colors.HexColor("#2F3640"),
        )

        story = [Paragraph(title, title_style), Spacer(1, 20)]

        for section in sections:
            story.append(HRFlowable(width="100%", color=colors.HexColor("#BDC3C7"), thickness=0.8))
            story.append(Spacer(1, 8))
            story.append(Paragraph(section["heading"], heading_style))
            story.append(Paragraph(section["content"], body_style))
            story.append(Spacer(1, 12))

        doc.build(story)

    # === POWERPOINT GENERATION ===
    elif file_type == "pptx":
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = "Generated automatically"
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

        for section in sections:
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)

            # Heading box
            left, top, width, height = Inches(0.8), Inches(0.8), Inches(8.4), Inches(1)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            p = title_frame.add_paragraph()
            p.text = section["heading"]
            p.font.bold = True
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(26, 188, 156)

            # Content box
            left, top, width, height = Inches(0.8), Inches(2), Inches(8.4), Inches(4.5)
            content_box = slide.shapes.add_textbox(left, top, width, height)
            tf = content_box.text_frame
            p = tf.add_paragraph()
            p.text = section["content"]
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(52, 73, 94)
            p.line_spacing = 1.3

        prs.save(str(file_path))

    else:
        raise ValueError("file_type must be either 'pdf' or 'pptx'")

    return str(file_path)


def clean_think_blocks(text: str) -> str:
    """
    Removes all <think>...</think> sections from the given text.
    """
    cleaned_text = re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL)
    return cleaned_text.strip()

def extract_assistant_response(text: str) -> str:
    """
    Extracts the Assistant's actual response text from a raw SmolVLM output.

    Example:
        Input:
            "🧠 Response: User: You are provided the following... Assistant: Yes, there is a graph..."
        Output:
            "Yes, there is a graph in the video. It shows the progress..."
    """
    # Remove leading emojis or response prefixes like "🧠 Response:"
    cleaned_text = re.sub(r"^🧠\s*Response:\s*", "", text.strip())

    # Use regex to find text after 'Assistant:'
    match = re.search(r"Assistant:\s*(.*)", cleaned_text, re.DOTALL)
    if match:
        # Clean trailing whitespace/newlines
        return match.group(1).strip()
    return cleaned_text.strip()